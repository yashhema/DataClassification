# src/extraction/content_extractor.py
"""
Universal Content Extraction Module for File-Based Connectors

This module provides comprehensive content extraction capabilities for all supported
file types, with memory-safe processing, error isolation, and configurable limits.

FIXES APPLIED: All critical bugs addressed
- Fixed infinite recursion bug in archive processing
- Added archive depth control and circular reference detection
- Corrected component type classification for archive members
- Added missing PyMuPDF table extraction
- Integrated error_handler and memory validation
- Fixed temp path generation with context IDs
- Implemented all missing file type processors
ASYNC CONVERSION: Converted all file I/O to async operations
"""

import os
import uuid
import aiofiles.os
import tempfile
import aiofiles
from typing import AsyncIterator, Dict, Any, Optional,List
from pathlib import Path
from datetime import datetime, timezone

# Core system imports
from core.models.models import ContentComponent, ContentExtractionConfig
from core.logging.system_logger import SystemLogger
from core.config.configuration_manager import SystemConfig
from core.errors import ErrorHandler
import aiofiles.os
import asyncio
import threading
# Add these imports to the top of content_extraction/content_extractor.py
import asyncio
try:
    import tikka
except ImportError:
    pass # Handled in __init__

# A lock to ensure the JVM is only started once per process
_tika_init_lock = threading.Lock()
_tika_initialized = False
class ContentExtractor:
    """Universal content extraction module for file-based connectors"""
    
    def __init__(self, extraction_config: ContentExtractionConfig, 
                 system_config: SystemConfig, logger: SystemLogger, error_handler: ErrorHandler):
        self.extraction_config = extraction_config
        self.system_config = system_config
        self.logger = logger
        self.error_handler = error_handler
        
        # Initialize extraction state
        self.processed_archive_paths = set()
        self.temp_files_created = []
        
        # Setup OCR if available
        self.ocr_reader = None
        self.ocr_available = False
        self._setup_ocr_reader()
        
        # Build extractor registry based on available libraries
        self.extractor_registry = self._build_extractor_registry()
        
        # Validate configuration
        self._validate_extraction_config()
        # UPDATED: Initialize the in-process Tika JVM based on the SYSTEM-LEVEL flag.
        global _tika_initialized
        
        # This is the master switch for the entire worker process.
        if self.system_config.connector.tika_fallback_enabled:
            with _tika_init_lock:
                if not _tika_initialized:
                    try:
                        import tikka
                        # This starts the JVM in-process via JPype.
                        tikka.init_tika() 
                        _tika_initialized = True
                        self.logger.info("Tika (JPype) has been initialized for the worker process.")
                    except Exception as e:
                        self.logger.error(f"System-level Tika fallback is enabled but failed to initialize. Fallback will be disabled for all datasources. Error: {e}")
                        _tika_initialized = False # Ensure we don't try again
    # =============================================================================
    # Main Extraction Interface
    # =============================================================================

    async def extract_from_file(self, file_path: str, object_id: str, 
                               job_id: int = 0, task_id: str = "0", datasource_id: str = "unknown") -> AsyncIterator[ContentComponent]:
        """
        Main extraction method - routes file to appropriate extractor
        """
        extraction_id = f"extract_{uuid.uuid4().hex[:8]}"
        
        try:
            # Pre-validation - check file size limits
            if not await aiofiles.os.path.exists(file_path):
                yield self._create_error_component(object_id, f"File not found: {file_path}")
                return
            
            
            # Memory validation before processing
            stats = (await aiofiles.os.stat(file_path)).st_size
            file_size_mb = stats / (1024 * 1024)
            if not self._validate_extraction_resources(file_size_mb):
                yield self._create_error_component(object_id, 
                    f"Insufficient memory for file processing (size: {file_size_mb:.2f}MB)")
                return
            
            
            
            if file_size_mb > self.extraction_config.limits.max_file_size_mb:
                yield self._create_oversized_component(object_id, {
                    "path": file_path,
                    "size_mb": file_size_mb,
                    "size_bytes": (await aiofiles.os.stat(file_path)).st_size
                })
                return
            
            # Detect file type and route to appropriate extractor
            file_type = self._detect_file_type(file_path)
            
            self.logger.info("Starting content extraction",
                           file_path=file_path,
                           object_id=object_id,
                           file_type=file_type,
                           file_size_mb=round(file_size_mb, 2),
                           extraction_id=extraction_id)
            
            # Route to appropriate extractor
            if file_type in ['zip', 'tar']:
                async for component in self._process_archive(file_path, object_id, file_type, 
                                                           job_id, task_id, datasource_id):
                    yield component
            elif file_type in self.extractor_registry:
                extractor = self.extractor_registry[file_type]
                async for component in extractor(file_path, object_id, job_id, task_id, datasource_id):
                    yield component
            else: # This block handles unsupported file types
                # UPDATED: Perform the two-level check
                global _tika_initialized
                use_tika = (
                    _tika_initialized and # Check if the system-level init was successful
                    self.extraction_config.features.tika_fallback_enabled # Check the data source-level flag
                )

                if use_tika:
                    self.logger.info(f"File type '{file_type}' not natively supported. Attempting Tika fallback.", object_id=object_id)
                    async for component in self._process_with_tika_fallback(file_path, object_id):
                        yield component
                else:
                    # Original logic if Tika is not enabled for this datasource or failed to initialize
                    yield self._create_unsupported_component(object_id, file_path, file_type)            
            self.logger.info("Content extraction completed",
                           extraction_id=extraction_id,
                           object_id=object_id)
            
        except Exception as e:
            error = self.error_handler.handle_error(
                e, f"extraction_system_error_{object_id}",
                operation="content_extraction",
                file_path=file_path,
                object_id=object_id
            )
            self.logger.error("Content extraction system error", error_id=error.error_id)
            yield self._create_error_component(object_id, f"System extraction error: {str(e)}")
        
        finally:
            await self._cleanup_temp_files_async(extraction_id)






    async def _process_with_tika_fallback(self, file_path: str, object_id: str) -> AsyncIterator[ContentComponent]:
        """
        Uses the in-process Tika (JPype) to extract content from unsupported file types.
        """
        try:
            # Get the current asyncio event loop to run the synchronous call in a thread.
            loop = asyncio.get_running_loop()
            
            # The tikka.parse() call is synchronous and can block for several seconds.
            # It MUST be run in an executor to avoid freezing the async worker.
            parsed = await loop.run_in_executor(
                None,       # Use the default thread pool executor
                tikka.parse, # The function to run
                file_path    # The argument to the function
            )
            
            # Safely get the content and metadata from Tika's response.
            content = parsed.get("content", "").strip() if parsed and "content" in parsed else ""
            
            if content:
                # If Tika successfully extracted text, create a standard text component.
                base_name = self._get_base_name(file_path)
                tika_metadata = parsed.get("metadata", {})
                
                # Add some of Tika's metadata for better context.
                extra_metadata = {
                    "tika_content_type": tika_metadata.get("Content-Type"),
                    "tika_producer": tika_metadata.get("producer")
                }
                
                # Use the existing helper to ensure component size limits are respected.
                component = self._create_text_component(
                    object_id=object_id,
                    component_id=f"{base_name}_tika_text_1",
                    file_path=file_path,
                    content=content,
                    extra_metadata=extra_metadata,
                    extraction_method="tika_fallback"
                )
                yield component
            else:
                # Log if Tika runs but finds no text content.
                self.logger.warning("Tika fallback ran but extracted no content.", object_id=object_id, file_path=file_path)

        except Exception as e:
            # If the Tika/JPype process fails, log the error and yield a
            # standard error component to ensure the failure is recorded.
            self.logger.error(f"Tika fallback processing failed for {file_path}", exc_info=True)
            yield self._create_error_component(object_id, f"Tika fallback failed: {str(e)}")




    # =============================================================================
    # Archive Processing - Fixed Recursion and Depth Control
    # =============================================================================

    async def _process_archive(self, archive_path: str, object_id: str, file_type: str,
                              job_id: int, task_id: str, datasource_id: str, 
                              current_depth: int = 0) -> AsyncIterator[ContentComponent]:
        """
        Process archive with depth control and circular reference detection
        """
        
        # Check depth limits
        if self.extraction_config.limits.max_archive_depth != -1:
            if current_depth >= self.extraction_config.limits.max_archive_depth:
                yield self._create_depth_exceeded_component(object_id, current_depth)
                return
        
        # Circular reference detection
        if archive_path in self.processed_archive_paths:
            self.logger.warning("Circular archive reference detected",
                               archive_path=archive_path)
            yield self._create_circular_reference_component(object_id, archive_path)
            return
        
        self.processed_archive_paths.add(archive_path)
        
        try:
            if file_type == 'zip':
                async for component in self._process_zip_archive(archive_path, object_id, 
                                                               job_id, task_id, datasource_id, current_depth):
                    yield component
            elif file_type == 'tar':
                async for component in self._process_tar_archive(archive_path, object_id,
                                                               job_id, task_id, datasource_id, current_depth):
                    yield component
        finally:
            self.processed_archive_paths.discard(archive_path)

    async def _process_zip_archive(self, archive_path: str, object_id: str,
                                 job_id: int, task_id: str, datasource_id: str,
                                 current_depth: int) -> AsyncIterator[ContentComponent]:
        """Process ZIP archive without recursion - Non-blocking version"""
        
        try:
            # Run the blocking extraction in a thread pool
            loop = asyncio.get_running_loop()
            components = await loop.run_in_executor(
                None, 
                self._blocking_zip_archive_processing,
                archive_path, object_id, job_id, task_id, datasource_id, current_depth
            )
            
            # Yield each component
            for component in components:
                yield component
                
        except Exception as e:
            error = self.error_handler.handle_error(
                e, f"zip_processing_{object_id}",
                operation="zip_archive_processing",
                archive_path=archive_path
            )
            self.logger.error("ZIP processing failed", error_id=error.error_id)
            yield self._create_error_component(object_id, f"ZIP processing failed: {str(e)}")

    def _blocking_zip_archive_processing(self, archive_path: str, object_id: str,
                                       job_id: int, task_id: str, datasource_id: str,
                                       current_depth: int) -> List[ContentComponent]:
        """Synchronous ZIP archive processing to run in thread pool"""
        components = []
        
        try:
            import zipfile
            
            with zipfile.ZipFile(archive_path, 'r') as zf:
                member_names = zf.namelist()
                processed_members = 0
                
                for member_name in member_names:
                    # Security validation
                    if not self._validate_archive_member_path(member_name):
                        self.logger.warning("Unsafe ZIP member path skipped",
                                           member_path=member_name)
                        continue
                    
                    # Member count limits
                    if processed_members >= self.extraction_config.limits.max_archive_members:
                        break
                    
                    # Skip directories
                    if member_name.endswith('/'):
                        continue
                    
                    try:
                        # Extract member to temp file
                        member_data = zf.read(member_name)
                        temp_member_path = self._create_temp_path(job_id, task_id, datasource_id, 
                                                                f"zip_member_{processed_members}")
                        
                        # Write file synchronously since we're in thread pool
                        with open(temp_member_path, 'wb') as temp_file:
                            temp_file.write(member_data)
                        self.temp_files_created.append(temp_member_path)
                        
                        # Process member
                        member_object_id = f"{object_id}/{member_name}"
                        member_file_type = self._detect_file_type(temp_member_path)
                        
                        # Check if member is also an archive (nested archives)
                        if member_file_type in ['zip', 'tar']:
                            # For nested archives, we'll need to handle this differently in the async context
                            # Store info for async processing later
                            nested_archive_component = ContentComponent(
                                object_id=member_object_id,
                                component_type="nested_archive_placeholder",
                                component_id=f"{self._get_base_name(temp_member_path)}_nested_archive_1",
                                parent_path=f"{archive_path}/{member_name}",
                                content=f"Nested archive: {member_file_type}",
                                original_size=len(member_data),
                                extracted_size=0,
                                is_truncated=False,
                                schema={"temp_path": temp_member_path, "file_type": member_file_type, "depth": current_depth + 1},
                                metadata={
                                    "archive_parent": archive_path,
                                    "member_name": member_name,
                                    "requires_async_processing": True
                                },
                                extraction_method="zip_nested_archive_marker",
                                is_archive_extraction=True
                            )
                            components.append(nested_archive_component)
                        else:
                            member_components = self._process_archive_member_blocking(temp_member_path, member_object_id,
                                                                                    member_file_type, archive_path, member_name,
                                                                                    job_id, task_id, datasource_id)
                            components.extend(member_components)
                        
                        processed_members += 1
                        
                    except Exception as member_error:
                        self.logger.warning("Failed to process ZIP member",
                                           member_name=member_name,
                                           error=str(member_error))
                        continue
            
            return components
                        
        except Exception as e:
            error_component = self._create_error_component(object_id, f"ZIP processing failed: {str(e)}")
            return [error_component]

    def _process_archive_member_blocking(self, member_path: str, member_object_id: str, 
                                       member_file_type: str, archive_path: str, member_name: str,
                                       job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Process archive member directly without recursion - blocking version"""
        
        components = []
        
        # Route directly to file type extractor (synchronous versions)
        if member_file_type in self.extractor_registry:
            # We need to call the blocking versions of extractors here
            if member_file_type == 'pdf':
                member_components = self._blocking_pdf_extraction(member_path, member_object_id, job_id, task_id, datasource_id)
            elif member_file_type == 'docx':
                member_components = self._blocking_docx_extraction(member_path, member_object_id, job_id, task_id, datasource_id)
            elif member_file_type == 'xlsx':
                member_components = self._blocking_xlsx_extraction(member_path, member_object_id, job_id, task_id, datasource_id)
            elif member_file_type == 'pptx':
                member_components = self._blocking_pptx_extraction(member_path, member_object_id, job_id, task_id, datasource_id)
            elif member_file_type == 'xml':
                member_components = self._blocking_xml_extraction(member_path, member_object_id, job_id, task_id, datasource_id)
            elif member_file_type == 'eml':
                member_components = self._blocking_eml_extraction(member_path, member_object_id, job_id, task_id, datasource_id)
            elif member_file_type == 'html':
                member_components = self._blocking_html_extraction(member_path, member_object_id, job_id, task_id, datasource_id)
            elif member_file_type == 'json':
                member_components = self._blocking_json_extraction(member_path, member_object_id, job_id, task_id, datasource_id)
            elif member_file_type == 'image':
                member_components = self._blocking_image_extraction(member_path, member_object_id, job_id, task_id, datasource_id)
            elif member_file_type == 'txt':
                # For txt files, simple synchronous read
                try:
                    with open(member_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text_content = f.read()
                    
                    text_component = self._create_text_component(
                        member_object_id, f"{self._get_base_name(member_path)}_txt_text_1", member_path,
                        text_content, {"encoding": "utf-8"}, "text_file_read"
                    )
                    member_components = [text_component]
                except Exception as txt_error:
                    error_component = self._create_error_component(member_object_id, f"Text file extraction failed: {str(txt_error)}")
                    member_components = [error_component]
            else:
                # Fallback for unhandled types
                member_components = []
            
            # Set component_type to 'archive_member' for all archive content
            for component in member_components:
                original_type = component.component_type
                component.component_type = "archive_member"
                component.is_archive_extraction = True
                component.parent_path = f"{archive_path}/{member_name}"
                
                # Preserve original type in metadata
                if not hasattr(component, 'metadata') or component.metadata is None:
                    component.metadata = {}
                component.metadata["original_component_type"] = original_type
                component.metadata["archive_parent"] = archive_path
                component.metadata["member_name"] = member_name
            
            components.extend(member_components)
        else:
            # Unsupported member type
            unsupported_component = ContentComponent(
                object_id=member_object_id,
                component_type="archive_member",
                component_id=f"{self._get_base_name(member_path)}_unsupported_1",
                parent_path=f"{archive_path}/{member_name}",
                content=f"Archive member type '{member_file_type}' not supported",
                original_size=self._get_file_size_sync(member_path),
                extracted_size=0,
                is_truncated=False,
                schema={},
                metadata={
                    "original_component_type": "unsupported_format",
                    "archive_parent": archive_path,
                    "member_name": member_name
                },
                extraction_method="archive_member_unsupported",
                is_archive_extraction=True
            )
            components.append(unsupported_component)
        
        return components

    async def _process_tar_archive(self, archive_path: str, object_id: str,
                                 job_id: int, task_id: str, datasource_id: str,
                                 current_depth: int) -> AsyncIterator[ContentComponent]:
        """Process TAR archive without recursion - Non-blocking version"""
        
        try:
            # Run the blocking extraction in a thread pool
            loop = asyncio.get_running_loop()
            components = await loop.run_in_executor(
                None, 
                self._blocking_tar_archive_processing,
                archive_path, object_id, job_id, task_id, datasource_id, current_depth
            )
            
            # Yield each component
            for component in components:
                yield component
                
        except Exception as e:
            error = self.error_handler.handle_error(
                e, f"tar_processing_{object_id}",
                operation="tar_archive_processing",
                archive_path=archive_path
            )
            self.logger.error("TAR processing failed", error_id=error.error_id)
            yield self._create_error_component(object_id, f"TAR processing failed: {str(e)}")

    def _blocking_tar_archive_processing(self, archive_path: str, object_id: str,
                                       job_id: int, task_id: str, datasource_id: str,
                                       current_depth: int) -> List[ContentComponent]:
        """Synchronous TAR archive processing to run in thread pool"""
        components = []
        
        try:
            import tarfile
            
            with tarfile.open(archive_path, 'r') as tf:
                members = tf.getmembers()
                processed_members = 0
                
                for member in members:
                    # Security validation
                    if not self._validate_archive_member_path(member.name):
                        self.logger.warning("Unsafe TAR member path skipped",
                                           member_path=member.name)
                        continue
                    
                    # Member count limits
                    if processed_members >= self.extraction_config.limits.max_archive_members:
                        break
                    
                    # Only process files
                    if not member.isfile():
                        continue
                    
                    try:
                        # Extract member content
                        extracted_file = tf.extractfile(member)
                        if extracted_file is None:
                            continue
                        
                        member_data = extracted_file.read()
                        temp_member_path = self._create_temp_path(job_id, task_id, datasource_id,
                                                                f"tar_member_{processed_members}")
                        
                        # Write file synchronously since we're in thread pool
                        with open(temp_member_path, 'wb') as temp_file:
                            temp_file.write(member_data)
                        self.temp_files_created.append(temp_member_path)
                        
                        # Process member
                        member_object_id = f"{object_id}/{member.name}"
                        member_file_type = self._detect_file_type(temp_member_path)
                        
                        # Check if member is also an archive (nested archives)
                        if member_file_type in ['zip', 'tar']:
                            # For nested archives, we'll need to handle this differently in the async context
                            # Store info for async processing later
                            nested_archive_component = ContentComponent(
                                object_id=member_object_id,
                                component_type="nested_archive_placeholder",
                                component_id=f"{self._get_base_name(temp_member_path)}_nested_archive_1",
                                parent_path=f"{archive_path}/{member.name}",
                                content=f"Nested archive: {member_file_type}",
                                original_size=len(member_data),
                                extracted_size=0,
                                is_truncated=False,
                                schema={"temp_path": temp_member_path, "file_type": member_file_type, "depth": current_depth + 1},
                                metadata={
                                    "archive_parent": archive_path,
                                    "member_name": member.name,
                                    "requires_async_processing": True
                                },
                                extraction_method="tar_nested_archive_marker",
                                is_archive_extraction=True
                            )
                            components.append(nested_archive_component)
                        else:
                            member_components = self._process_archive_member_blocking(temp_member_path, member_object_id,
                                                                                    member_file_type, archive_path, member.name,
                                                                                    job_id, task_id, datasource_id)
                            components.extend(member_components)
                        
                        processed_members += 1
                        
                    except Exception as member_error:
                        self.logger.warning("Failed to process TAR member",
                                           member_name=member.name,
                                           error=str(member_error))
                        continue
            
            return components
                        
        except Exception as e:
            error_component = self._create_error_component(object_id, f"TAR processing failed: {str(e)}")
            return [error_component]

    # =============================================================================
    # Memory and Resource Validation
    # =============================================================================

    def _validate_extraction_resources(self, file_size_mb: float) -> bool:
        """
        Check if file can be processed within resource limits
        TODO: Replace with ResourceCoordinator integration for real memory monitoring
        """
        try:
            # Simple memory check - could be enhanced with actual memory monitoring
            total_memory_limit = self.system_config.system.total_process_memory_limit_mb
            safe_memory_limit = total_memory_limit * 0.7  # Reserve 30% for overhead
            
            if file_size_mb > safe_memory_limit:
                self.logger.warning("File size exceeds safe memory limit",
                                   file_size_mb=file_size_mb,
                                   safe_limit_mb=safe_memory_limit)
                return False
            
            return True
            
        except Exception as e:
            self.logger.warning("Memory validation failed", error=str(e))
            return True  # Default to allowing processing

    # =============================================================================
    # PDF Processing - Added Missing Table Extraction
    # =============================================================================
    async def extract_pdf_components(self, file_path: str, object_id: str,
                                   job_id: int, task_id: str, datasource_id: str) -> AsyncIterator[ContentComponent]:
        """Complete PDF content extraction with PyMuPDF table support - Non-blocking version"""
        
        try:
            # Run the blocking extraction in a thread pool
            loop = asyncio.get_running_loop()
            components = await loop.run_in_executor(
                None, 
                self._blocking_pdf_extraction,
                file_path, object_id, job_id, task_id, datasource_id
            )
            
            # Yield each component
            for component in components:
                yield component
                
        except Exception as e:
            yield self._create_error_component(object_id, f"PDF processing system error: {str(e)}")

    def _blocking_pdf_extraction(self, file_path: str, object_id: str, 
                                job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Synchronous PDF extraction to run in thread pool"""
        components = []
        
        try:
            import pymupdf as fitz
            doc = fitz.open(file_path)
            
            try:
                base_name = self._get_base_name(file_path)
                
                # 1. Extract text content
                text_content = ""
                for page in doc:
                    text_content += page.get_text()
                
                if text_content.strip():
                    text_component = self._create_text_component(
                        object_id, f"{base_name}_pdf_text_1", file_path,
                        text_content, {"page_count": len(doc)}, "pymupdf_text"
                    )
                    components.append(text_component)
                
                # 2. Extract tables (if enabled)
                if self.extraction_config.features.extract_tables:
                    table_components = self._extract_pdf_tables_pymupdf_blocking(file_path, object_id, doc)
                    components.extend(table_components)
                
                # 3. Extract images (if enabled and OCR available)
                if (self.extraction_config.features.extract_pictures and 
                    self._is_ocr_enabled()):
                    image_components = self._extract_pdf_images_blocking(file_path, object_id, doc, job_id, task_id, datasource_id)
                    components.extend(image_components)
                
            finally:
                doc.close()
                
            return components
            
        except ImportError:
            error_component = self._create_error_component(object_id, "PyMuPDF library not available")
            return [error_component]
        except Exception as e:
            error_component = self._create_error_component(object_id, f"PDF extraction failed: {str(e)}")
            return [error_component]

    def _extract_pdf_tables_pymupdf_blocking(self, file_path: str, object_id: str, doc) -> List[ContentComponent]:
        """PyMuPDF table extraction as blocking function for thread pool"""
        
        components = []
        table_index = 1
        base_name = self._get_base_name(file_path)
        tables_found = False
        
        try:
            # PRIMARY: PyMuPDF table extraction
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                table_finder = page.find_tables()
                
                for table in table_finder.tables:
                    if table_index > self.extraction_config.limits.max_document_table_rows:
                        break
                    
                    # Convert table to pandas DataFrame
                    df = table.to_pandas()
                    table_data = self._convert_pymupdf_table(df, table_index)
                    
                    component = ContentComponent(
                        object_id=object_id,
                        component_type="table",
                        component_id=f"{base_name}_pdf_table_{table_index}",
                        parent_path=file_path,
                        content="",  # Empty content for tables
                        original_size=len(str(df)),
                        extracted_size=0,
                        is_truncated=False,
                        schema=table_data,
                        metadata=self._get_standard_metadata("table", {"page_number": page_num + 1}),
                        extraction_method="pymupdf_tables"
                    )
                    components.append(component)
                    
                    table_index += 1
                    tables_found = True
                    
        except Exception as pymupdf_error:
            self.logger.warning("PyMuPDF table extraction failed, trying Camelot",
                               file_path=file_path, error=str(pymupdf_error))
        
        # FALLBACK: Camelot if PyMuPDF failed
        if not tables_found:
            camelot_components = self._extract_pdf_tables_camelot_blocking(file_path, object_id, table_index)
            components.extend(camelot_components)
            
        return components

    def _extract_pdf_tables_camelot_blocking(self, file_path: str, object_id: str, table_index: int) -> List[ContentComponent]:
        """Camelot fallback for PDF table extraction - blocking version"""
        
        components = []
        base_name = self._get_base_name(file_path)
        
        try:
            import camelot
            tables = camelot.read_pdf(file_path, pages='all')
            
            for table in tables:
                if table_index > self.extraction_config.limits.max_document_table_rows:
                    break
                
                table_data = self._convert_camelot_table(table)
                
                component = ContentComponent(
                    object_id=object_id,
                    component_type="table",
                    component_id=f"{base_name}_pdf_table_{table_index}",
                    parent_path=file_path,
                    content="",
                    original_size=len(table_data["serialized_rows"]),
                    extracted_size=0,
                    is_truncated=False,
                    schema=table_data,
                    metadata=self._get_standard_metadata("table", {"extraction_library": "camelot"}),
                    extraction_method="camelot"
                )
                components.append(component)
                table_index += 1
                
        except Exception as camelot_error:
            # Try Tabula as final fallback
            try:
                import tabula
                tables = tabula.read_pdf(file_path, pages='all', lattice=True)
                
                for table_df in tables:
                    table_data = self._convert_tabula_table(table_df)
                    
                    component = ContentComponent(
                        object_id=object_id,
                        component_type="table",
                        component_id=f"{base_name}_pdf_table_{table_index}",
                        parent_path=file_path,
                        content="",
                        original_size=len(table_data["serialized_rows"]),
                        extracted_size=0,
                        is_truncated=False,
                        schema=table_data,
                        metadata=self._get_standard_metadata("table", {"extraction_library": "tabula"}),
                        extraction_method="tabula_fallback"
                    )
                    components.append(component)
                    table_index += 1
                    
            except Exception as tabula_error:
                self.logger.error("All PDF table extraction methods failed",
                                 camelot_error=str(camelot_error),
                                 tabula_error=str(tabula_error))
        
        return components

    def _extract_pdf_images_blocking(self, file_path: str, object_id: str, doc, 
                                   job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Extract and OCR images from PDF - blocking version"""
        import pymupdf as fitz
        
        components = []
        
        if not self.ocr_available:
            return components
        
        base_name = self._get_base_name(file_path)
        image_index = 1
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    
                    if pix.n - pix.alpha < 4:  # Skip CMYK
                        temp_img_path = self._create_temp_path(job_id, task_id, datasource_id, 
                                                             f"pdf_img_{page_num}_{img_index}")
                        pix.save(temp_img_path)
                        self.temp_files_created.append(temp_img_path)
                        
                        # Use synchronous OCR since we're already in thread pool
                        ocr_text = self._extract_text_from_image(temp_img_path)
                        if ocr_text.strip():
                            # Get file size synchronously since we're in thread pool
                            import os
                            file_size = os.path.getsize(temp_img_path)
                            
                            component = ContentComponent(
                                object_id=object_id,
                                component_type="image_ocr",
                                component_id=f"{base_name}_pdf_image_ocr_{image_index}",
                                parent_path=file_path,
                                content=ocr_text,
                                original_size=file_size,
                                extracted_size=len(ocr_text),
                                is_truncated=False,
                                schema={},
                                metadata=self._get_standard_metadata("image_ocr", {"page_number": page_num + 1}),
                                extraction_method="pdf_image_ocr"
                            )
                            components.append(component)
                        
                        image_index += 1
                    
                    pix = None  # Free memory
                    
                except Exception as img_error:
                    self.logger.warning("Failed to process PDF image",
                                       page_num=page_num, error=str(img_error))
                    continue
        
        return components


    # =============================================================================
    # Word Document Processing - Added Missing Image Extraction
    # =============================================================================
    async def extract_docx_components(self, file_path: str, object_id: str,
                                    job_id: int, task_id: str, datasource_id: str) -> AsyncIterator[ContentComponent]:
        """Complete Word document extraction - Non-blocking version"""
        
        try:
            # Run the blocking extraction in a thread pool
            loop = asyncio.get_running_loop()
            components = await loop.run_in_executor(
                None, 
                self._blocking_docx_extraction,
                file_path, object_id, job_id, task_id, datasource_id
            )
            
            # Yield each component
            for component in components:
                yield component
                
        except Exception as e:
            yield self._create_error_component(object_id, f"Word document processing system error: {str(e)}")

    def _blocking_docx_extraction(self, file_path: str, object_id: str, 
                                 job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Synchronous Word document extraction to run in thread pool"""
        components = []
        
        try:
            from docx import Document
            
            doc = Document(file_path)
            base_name = self._get_base_name(file_path)
            
            # 1. Extract text content
            text_content = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            
            if text_content.strip():
                text_component = self._create_text_component(
                    object_id, f"{base_name}_docx_text_1", file_path,
                    text_content, {"paragraph_count": len(doc.paragraphs)}, "python_docx_text"
                )
                components.append(text_component)
            
            # 2. Extract tables (if enabled)
            if self.extraction_config.features.extract_tables:
                for table_index, table in enumerate(doc.tables, 1):
                    table_data = self._convert_docx_table(table)
                    
                    if table_data["rows"]:
                        component = ContentComponent(
                            object_id=object_id,
                            component_type="table",
                            component_id=f"{base_name}_docx_table_{table_index}",
                            parent_path=file_path,
                            content="",
                            original_size=len(table_data["serialized_rows"]),
                            extracted_size=0,
                            is_truncated=False,
                            schema=table_data,
                            metadata=self._get_standard_metadata("table", {"table_index": table_index}),
                            extraction_method="python_docx_tables"
                        )
                        components.append(component)
            
            # 3. Extract images (if enabled)
            if (self.extraction_config.features.extract_pictures and self._is_ocr_enabled()):
                image_components = self._extract_docx_images_blocking(file_path, object_id, doc, job_id, task_id, datasource_id)
                components.extend(image_components)
            
            return components
            
        except ImportError:
            error_component = self._create_error_component(object_id, "python-docx library not available")
            return [error_component]
        except Exception as e:
            error_component = self._create_error_component(object_id, f"Word document extraction failed: {str(e)}")
            return [error_component]

    def _extract_docx_images_blocking(self, file_path: str, object_id: str, doc, 
                                    job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Extract images from Word document - blocking version"""
        
        components = []
        
        if not self.ocr_available:
            return components
        
        try:
            import zipfile
            import os
            base_name = self._get_base_name(file_path)
            image_index = 1
            
            # Word documents are ZIP archives - extract images from word/media/
            with zipfile.ZipFile(file_path, 'r') as docx_zip:
                image_files = [f for f in docx_zip.namelist() if f.startswith('word/media/')]
                
                for img_file in image_files:
                    try:
                        img_data = docx_zip.read(img_file)
                        temp_img_path = self._create_temp_path(job_id, task_id, datasource_id, 
                                                             f"docx_img_{image_index}")
                        
                        # Write file synchronously since we're in thread pool
                        with open(temp_img_path, 'wb') as f:
                            f.write(img_data)
                        self.temp_files_created.append(temp_img_path)
                        
                        # Use synchronous OCR since we're already in thread pool
                        ocr_text = self._extract_text_from_image(temp_img_path)
                        if ocr_text.strip():
                            component = ContentComponent(
                                object_id=object_id,
                                component_type="image_ocr",
                                component_id=f"{base_name}_docx_image_ocr_{image_index}",
                                parent_path=file_path,
                                content=ocr_text,
                                original_size=len(img_data),
                                extracted_size=len(ocr_text),
                                is_truncated=False,
                                schema={},
                                metadata=self._get_standard_metadata("image_ocr", {"source_file": img_file}),
                                extraction_method="docx_image_ocr"
                            )
                            components.append(component)
                        
                        image_index += 1
                        
                    except Exception as img_error:
                        self.logger.warning("Failed to process Word image",
                                           img_file=img_file, error=str(img_error))
                        continue
            
            return components
                        
        except Exception as e:
            self.logger.warning("Word image extraction failed", error=str(e))
            return components

    # =============================================================================
    # Excel Processing - Added Text Extraction
    # =============================================================================
    async def extract_xlsx_components(self, file_path: str, object_id: str,
                                    job_id: int, task_id: str, datasource_id: str) -> AsyncIterator[ContentComponent]:
        """Complete Excel workbook extraction with text and tables - Non-blocking version"""
        
        try:
            # Run the blocking extraction in a thread pool
            loop = asyncio.get_running_loop()
            components = await loop.run_in_executor(
                None, 
                self._blocking_xlsx_extraction,
                file_path, object_id, job_id, task_id, datasource_id
            )
            
            # Yield each component
            for component in components:
                yield component
                
        except Exception as e:
            yield self._create_error_component(object_id, f"Excel workbook processing system error: {str(e)}")

    def _blocking_xlsx_extraction(self, file_path: str, object_id: str, 
                                 job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Synchronous Excel workbook extraction to run in thread pool"""
        components = []
        
        try:
            import pandas as pd
            
            base_name = self._get_base_name(file_path)
            excel_file = pd.ExcelFile(file_path)
            
            # Extract text content from all cells
            all_text_content = []
            
            for sheet_index, sheet_name in enumerate(excel_file.sheet_names, 1):
                try:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    
                    if df.empty:
                        continue
                    
                    # Extract text from all cells
                    sheet_text = []
                    for column in df.columns:
                        column_text = df[column].astype(str).tolist()
                        sheet_text.extend([text for text in column_text if text and text != 'nan'])
                    
                    if sheet_text:
                        all_text_content.append(f"Sheet {sheet_name}: " + " ".join(sheet_text))
                    
                    # Extract table structure
                    table_data = self._convert_excel_sheet(df, sheet_name)
                    
                    component = ContentComponent(
                        object_id=object_id,
                        component_type="table",
                        component_id=f"{base_name}_xlsx_sheet_{sheet_index}",
                        parent_path=file_path,
                        content="",
                        original_size=len(table_data["serialized_rows"]),
                        extracted_size=0,
                        is_truncated=False,
                        schema=table_data,
                        metadata=self._get_standard_metadata("table", {"sheet_name": sheet_name}),
                        extraction_method="pandas_excel"
                    )
                    components.append(component)
                    
                except Exception as sheet_error:
                    self.logger.warning("Failed to process Excel sheet",
                                       sheet_name=sheet_name, error=str(sheet_error))
                    continue
            
            # Yield combined text content from all sheets
            if all_text_content:
                combined_text = '\n\n'.join(all_text_content)
                text_component = self._create_text_component(
                    object_id, f"{base_name}_xlsx_text_1", file_path,
                    combined_text, {"sheet_count": len(excel_file.sheet_names)}, "excel_cell_text"
                )
                components.append(text_component)
            
            # Extract embedded images (if enabled)
            if (self.extraction_config.features.extract_pictures and self._is_ocr_enabled()):
                image_components = self._extract_xlsx_images_blocking(file_path, object_id, job_id, task_id, datasource_id)
                components.extend(image_components)
            
            return components
            
        except ImportError:
            error_component = self._create_error_component(object_id, "pandas/openpyxl libraries not available")
            return [error_component]
        except Exception as e:
            error_component = self._create_error_component(object_id, f"Excel extraction failed: {str(e)}")
            return [error_component]

    def _extract_xlsx_images_blocking(self, file_path: str, object_id: str, 
                                    job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Extract images from Excel file - blocking version"""
        
        components = []
        
        if not self.ocr_available:
            return components
        
        try:
            import zipfile
            import os
            base_name = self._get_base_name(file_path)
            image_index = 1
            
            # Excel files are ZIP archives - extract from xl/media/
            with zipfile.ZipFile(file_path, 'r') as xlsx_zip:
                media_files = [f for f in xlsx_zip.namelist() if f.startswith('xl/media/')]
                
                for media_file in media_files:
                    try:
                        img_data = xlsx_zip.read(media_file)
                        temp_img_path = self._create_temp_path(job_id, task_id, datasource_id, 
                                                             f"xlsx_img_{image_index}")
                        
                        # Write file synchronously since we're in thread pool
                        with open(temp_img_path, 'wb') as f:
                            f.write(img_data)
                        self.temp_files_created.append(temp_img_path)
                        
                        # Use synchronous OCR since we're already in thread pool
                        ocr_text = self._extract_text_from_image(temp_img_path)
                        if ocr_text.strip():
                            component = ContentComponent(
                                object_id=object_id,
                                component_type="image_ocr",
                                component_id=f"{base_name}_xlsx_image_ocr_{image_index}",
                                parent_path=file_path,
                                content=ocr_text,
                                original_size=len(img_data),
                                extracted_size=len(ocr_text),
                                is_truncated=False,
                                schema={},
                                metadata=self._get_standard_metadata("image_ocr", {"source_file": media_file}),
                                extraction_method="xlsx_image_ocr"
                            )
                            components.append(component)
                        
                        image_index += 1
                        
                    except Exception as img_error:
                        self.logger.warning("Failed to process Excel image",
                                           media_file=media_file, error=str(img_error))
                        continue
            
            return components
                        
        except Exception as e:
            self.logger.warning("Excel image extraction failed", error=str(e))
            return components

    # =============================================================================
    # PowerPoint Processing - Added Missing Implementation
    # =============================================================================

    async def extract_pptx_components(self, file_path: str, object_id: str,
                                    job_id: int, task_id: str, datasource_id: str) -> AsyncIterator[ContentComponent]:
        """Complete PowerPoint presentation extraction - Non-blocking version"""
        
        try:
            # Run the blocking extraction in a thread pool
            loop = asyncio.get_running_loop()
            components = await loop.run_in_executor(
                None, 
                self._blocking_pptx_extraction,
                file_path, object_id, job_id, task_id, datasource_id
            )
            
            # Yield each component
            for component in components:
                yield component
                
        except Exception as e:
            yield self._create_error_component(object_id, f"PowerPoint presentation processing system error: {str(e)}")

    def _blocking_pptx_extraction(self, file_path: str, object_id: str, 
                                 job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Synchronous PowerPoint presentation extraction to run in thread pool"""
        components = []
        
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            base_name = self._get_base_name(file_path)
            
            # Extract text content from all slides
            all_text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                if slide_text:
                    all_text_content.append(f"Slide {slide_num}: " + " ".join(slide_text))
            
            # Yield combined text content
            if all_text_content:
                combined_text = '\n\n'.join(all_text_content)
                text_component = self._create_text_component(
                    object_id, f"{base_name}_pptx_text_1", file_path,
                    combined_text, {"slide_count": len(prs.slides)}, "python_pptx_text"
                )
                components.append(text_component)
            
            # Extract tables (if enabled)
            if self.extraction_config.features.extract_tables:
                table_index = 1
                for slide_num, slide in enumerate(prs.slides, 1):
                    for shape in slide.shapes:
                        if shape.has_table:
                            table_data = self._convert_pptx_table(shape.table)
                            
                            component = ContentComponent(
                                object_id=object_id,
                                component_type="table",
                                component_id=f"{base_name}_pptx_table_{table_index}",
                                parent_path=file_path,
                                content="",
                                original_size=len(table_data["serialized_rows"]),
                                extracted_size=0,
                                is_truncated=False,
                                schema=table_data,
                                metadata=self._get_standard_metadata("table", {"slide_number": slide_num}),
                                extraction_method="python_pptx_tables"
                            )
                            components.append(component)
                            table_index += 1
            
            # Extract images (if enabled)
            if (self.extraction_config.features.extract_pictures and self._is_ocr_enabled()):
                image_components = self._extract_pptx_images_blocking(file_path, object_id, prs, job_id, task_id, datasource_id)
                components.extend(image_components)
            
            return components
            
        except ImportError:
            error_component = self._create_error_component(object_id, "python-pptx library not available")
            return [error_component]
        except Exception as e:
            error_component = self._create_error_component(object_id, f"PowerPoint extraction failed: {str(e)}")
            return [error_component]

    def _extract_pptx_images_blocking(self, file_path: str, object_id: str, prs, 
                                    job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Extract images from PowerPoint presentation - blocking version"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        components = []
        
        if not self.ocr_available:
            return components
        
        base_name = self._get_base_name(file_path)
        image_index = 1
        
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        
                        temp_img_path = self._create_temp_path(job_id, task_id, datasource_id, 
                                                             f"pptx_img_{image_index}")
                        
                        # Write file synchronously since we're in thread pool
                        with open(temp_img_path, 'wb') as f:
                            f.write(image_bytes)
                        self.temp_files_created.append(temp_img_path)
                        
                        # Use synchronous OCR since we're already in thread pool
                        ocr_text = self._extract_text_from_image(temp_img_path)
                        if ocr_text.strip():
                            component = ContentComponent(
                                object_id=object_id,
                                component_type="image_ocr",
                                component_id=f"{base_name}_pptx_image_ocr_{image_index}",
                                parent_path=file_path,
                                content=ocr_text,
                                original_size=len(image_bytes),
                                extracted_size=len(ocr_text),
                                is_truncated=False,
                                schema={},
                                metadata=self._get_standard_metadata("image_ocr", {"slide_number": slide_num}),
                                extraction_method="pptx_image_ocr"
                            )
                            components.append(component)
                        
                        image_index += 1
                        
                    except Exception as img_error:
                        self.logger.warning("Failed to process PowerPoint image",
                                           slide_num=slide_num, error=str(img_error))
                        continue
        
        return components

    # =============================================================================
    # XML Processing - Complete Implementation
    # =============================================================================

    async def extract_xml_components(self, file_path: str, object_id: str,
                                   job_id: int, task_id: str, datasource_id: str) -> AsyncIterator[ContentComponent]:
        """Complete XML file extraction"""
        
        try:
            import xml.etree.ElementTree as ET
            
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                xml_content = await f.read()
            
            # Parse XML structure
            root = ET.fromstring(xml_content)
            base_name = self._get_base_name(file_path)
            
            # Extract text content from all elements
            all_text = []
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    all_text.append(elem.text.strip())
            
            if all_text:
                combined_text = ' '.join(all_text)
                text_component = self._create_text_component(
                    object_id, f"{base_name}_xml_text_1", file_path,
                    combined_text, {"root_tag": root.tag, "element_count": len(list(root.iter()))}, "xml_parsing"
                )
                yield text_component
            
            # Extract structured data if XML contains tabular data
            if self.extraction_config.features.extract_tables:
                table_data = self._extract_xml_tables(root)
                if table_data:
                    yield ContentComponent(
                        object_id=object_id,
                        component_type="table",
                        component_id=f"{base_name}_xml_table_1",
                        parent_path=file_path,
                        content="",
                        original_size=len(xml_content),
                        extracted_size=0,
                        is_truncated=False,
                        schema=table_data,
                        metadata=self._get_standard_metadata("table", {"source": "xml_elements"}),
                        extraction_method="xml_element_parsing"
                    )
            
        except Exception as e:
            yield self._create_error_component(object_id, f"XML extraction failed: {str(e)}")

    def _extract_xml_tables(self, root) -> Optional[Dict[str, Any]]:
        """Extract tabular data from XML elements"""
        
        # Look for repeating elements that could be table rows
        children = list(root)
        if len(children) < 2:
            return None
        
        # Check if children have similar structure (potential table rows)
        first_child = children[0]
        first_child_tags = [child.tag for child in first_child]
        
        # Verify other children have similar structure
        similar_structure = True
        for child in children[1:]:
            child_tags = [subchild.tag for subchild in child]
            if child_tags != first_child_tags:
                similar_structure = False
                break
        
        if not similar_structure or not first_child_tags:
            return None
        
        # Extract as table
        headers = first_child_tags
        rows = []
        
        for child in children:
            row_data = {}
            for subchild in child:
                row_data[subchild.tag] = subchild.text or ""
            rows.append(row_data)
        
        return {
            "headers": headers,
            "rows": rows,
            "row_count": len(rows),
            "column_count": len(headers),
            "serialized_rows": "\n".join(["|".join(row_data.values()) for row_data in rows])
        }

    # =============================================================================
    # Email Processing - Complete Implementation  
    # =============================================================================
    async def extract_eml_components(self, file_path: str, object_id: str,
                                   job_id: int, task_id: str, datasource_id: str) -> AsyncIterator[ContentComponent]:
        """Complete email file extraction - Non-blocking version"""
        
        try:
            # Run the blocking extraction in a thread pool
            loop = asyncio.get_running_loop()
            components = await loop.run_in_executor(
                None, 
                self._blocking_eml_extraction,
                file_path, object_id, job_id, task_id, datasource_id
            )
            
            # Yield each component
            for component in components:
                yield component
                
        except Exception as e:
            yield self._create_error_component(object_id, f"Email file processing system error: {str(e)}")

    def _blocking_eml_extraction(self, file_path: str, object_id: str, 
                                job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Synchronous email file extraction to run in thread pool"""
        components = []
        
        try:
            import email
            from email.policy import default
            
            # Read email data synchronously since we're in thread pool
            with open(file_path, 'rb') as f:
                msg_data = f.read()
                msg = email.message_from_bytes(msg_data, policy=default)
            
            base_name = self._get_base_name(file_path)
            
            # Extract email text content
            email_text = []
            
            # Get email headers
            headers_text = f"From: {msg['From']}\nTo: {msg['To']}\nSubject: {msg['Subject']}\nDate: {msg['Date']}"
            email_text.append(headers_text)
            
            # Get email body
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        email_text.append(part.get_content())
                    elif part.get_content_type() == "text/html":
                        # Parse HTML body for tables if enabled
                        if self.extraction_config.features.extract_tables:
                            table_components = self._extract_html_tables_from_content_blocking(
                                part.get_content(), object_id, f"{base_name}_email_html"
                            )
                            components.extend(table_components)
            else:
                email_text.append(msg.get_content())
            
            # Yield combined email text
            if email_text:
                combined_text = '\n\n'.join(email_text)
                text_component = self._create_text_component(
                    object_id, f"{base_name}_eml_text_1", file_path,
                    combined_text, {"from": msg['From'], "subject": msg['Subject']}, "email_parsing"
                )
                components.append(text_component)
            
            return components
            
        except Exception as e:
            error_component = self._create_error_component(object_id, f"Email extraction failed: {str(e)}")
            return [error_component]

    def _extract_html_tables_from_content_blocking(self, html_content: str, object_id: str, base_id: str) -> List[ContentComponent]:
        """Extract tables from HTML content string - blocking version"""
        
        components = []
        
        try:
            from bs4 import BeautifulSoup
            
            soup = BeautifulSoup(html_content, 'html.parser')
            tables = soup.find_all('table')
            
            for table_index, table in enumerate(tables, 1):
                table_data = self._convert_html_table(table)
                
                if table_data["rows"]:
                    component = ContentComponent(
                        object_id=object_id,
                        component_type="table",
                        component_id=f"{base_id}_table_{table_index}",
                        parent_path="",
                        content="",
                        original_size=len(table_data["serialized_rows"]),
                        extracted_size=0,
                        is_truncated=False,
                        schema=table_data,
                        metadata=self._get_standard_metadata("table", {"source": "html_email"}),
                        extraction_method="email_html_tables"
                    )
                    components.append(component)
            
            return components
                    
        except Exception as e:
            self.logger.warning("HTML table extraction from email failed", error=str(e))
            return components


    # =============================================================================
    # HTML Processing - Added Image Processing
    # =============================================================================
    async def extract_html_components(self, file_path: str, object_id: str,
                                    job_id: int, task_id: str, datasource_id: str) -> AsyncIterator[ContentComponent]:
        """Complete HTML document extraction with image processing - Non-blocking version"""
        
        try:
            # Run the blocking extraction in a thread pool
            loop = asyncio.get_running_loop()
            components = await loop.run_in_executor(
                None, 
                self._blocking_html_extraction,
                file_path, object_id, job_id, task_id, datasource_id
            )
            
            # Yield each component
            for component in components:
                yield component
                
        except Exception as e:
            yield self._create_error_component(object_id, f"HTML document processing system error: {str(e)}")

    def _blocking_html_extraction(self, file_path: str, object_id: str, 
                                 job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Synchronous HTML document extraction to run in thread pool"""
        components = []
        
        try:
            from bs4 import BeautifulSoup
            
            # Read HTML content synchronously since we're in thread pool
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            base_name = self._get_base_name(file_path)
            
            # 1. Extract text content
            text_content = soup.get_text()
            
            if text_content.strip():
                text_component = self._create_text_component(
                    object_id, f"{base_name}_html_text_1", file_path,
                    text_content, {"title": soup.title.string if soup.title else ""}, "beautifulsoup_text"
                )
                components.append(text_component)
            
            # 2. Extract tables (if enabled)
            if self.extraction_config.features.extract_tables:
                tables = soup.find_all('table')
                
                for table_index, table in enumerate(tables, 1):
                    table_data = self._convert_html_table(table)
                    
                    if table_data["rows"]:
                        component = ContentComponent(
                            object_id=object_id,
                            component_type="table",
                            component_id=f"{base_name}_html_table_{table_index}",
                            parent_path=file_path,
                            content="",
                            original_size=len(table_data["serialized_rows"]),
                            extracted_size=0,
                            is_truncated=False,
                            schema=table_data,
                            metadata=self._get_standard_metadata("table", {"table_index": table_index}),
                            extraction_method="beautifulsoup_tables"
                        )
                        components.append(component)
            
            # 3. Extract image references (if enabled)
            if (self.extraction_config.features.extract_pictures and self._is_ocr_enabled()):
                image_components = self._extract_html_images_blocking(file_path, object_id, soup, job_id, task_id, datasource_id)
                components.extend(image_components)
            
            return components
            
        except ImportError:
            error_component = self._create_error_component(object_id, "BeautifulSoup library not available")
            return [error_component]
        except Exception as e:
            error_component = self._create_error_component(object_id, f"HTML extraction failed: {str(e)}")
            return [error_component]

    def _extract_html_images_blocking(self, file_path: str, object_id: str, soup, 
                                    job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Extract and OCR images referenced in HTML - blocking version"""
        
        components = []
        
        if not self.ocr_available:
            return components
        
        images = soup.find_all('img')
        image_index = 1
        base_name = self._get_base_name(file_path)
        
        for img in images:
            src = img.get('src')
            if not src:
                continue
            
            try:
                # Resolve image path relative to HTML file
                import os
                html_dir = os.path.dirname(file_path)
                image_path = os.path.join(html_dir, src)
                
                if os.path.exists(image_path):
                    # Use synchronous OCR since we're already in thread pool
                    ocr_text = self._extract_text_from_image(image_path)
                    
                    if ocr_text.strip():
                        file_size = os.path.getsize(image_path)
                        component = ContentComponent(
                            object_id=object_id,
                            component_type="image_ocr", 
                            component_id=f"{base_name}_html_image_ocr_{image_index}",
                            parent_path=file_path,
                            content=ocr_text,
                            original_size=file_size,
                            extracted_size=len(ocr_text),
                            is_truncated=False,
                            schema={},
                            metadata=self._get_standard_metadata("image_ocr", {"src": src, "alt": img.get('alt', '')}),
                            extraction_method="html_image_ocr"
                        )
                        components.append(component)
                    
                    image_index += 1
                    
            except Exception as img_error:
                self.logger.warning("Failed to process HTML image",
                                   src=src, error=str(img_error))
                continue
        
        return components

    # =============================================================================
    # Simple File Type Extractors - Signatures Updated for Context IDs
    # =============================================================================

    async def extract_txt_components(self, file_path: str, object_id: str,
                                   job_id: int, task_id: str, datasource_id: str) -> AsyncIterator[ContentComponent]:
        """Extract text file content"""
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = await f.read()
            
            text_component = self._create_text_component(
                object_id, f"{self._get_base_name(file_path)}_txt_text_1", file_path,
                text_content, {"encoding": "utf-8"}, "text_file_read"
            )
            yield text_component
            
        except Exception as e:
            yield self._create_error_component(object_id, f"Text file extraction failed: {str(e)}")

    async def extract_json_components(self, file_path: str, object_id: str,
                                    job_id: int, task_id: str, datasource_id: str) -> AsyncIterator[ContentComponent]:
        """Extract JSON file content - Non-blocking version"""
        
        try:
            # Run the blocking extraction in a thread pool
            loop = asyncio.get_running_loop()
            components = await loop.run_in_executor(
                None, 
                self._blocking_json_extraction,
                file_path, object_id, job_id, task_id, datasource_id
            )
            
            # Yield each component
            for component in components:
                yield component
                
        except Exception as e:
            yield self._create_error_component(object_id, f"JSON file processing system error: {str(e)}")

    def _blocking_json_extraction(self, file_path: str, object_id: str, 
                                 job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Synchronous JSON file extraction to run in thread pool"""
        components = []
        
        try:
            import json
            
            # Read JSON content synchronously since we're in thread pool
            with open(file_path, 'r', encoding='utf-8') as f:
                json_text = f.read()
                json_data = json.loads(json_text)
            
            formatted_json = json.dumps(json_data, indent=2, ensure_ascii=False)
            
            text_component = self._create_text_component(
                object_id, f"{self._get_base_name(file_path)}_json_text_1", file_path,
                formatted_json, {"json_type": type(json_data).__name__}, "json_parsing"
            )
            components.append(text_component)
            
            return components
            
        except Exception as e:
            error_component = self._create_error_component(object_id, f"JSON extraction failed: {str(e)}")
            return [error_component]

    async def extract_image_components(self, file_path: str, object_id: str,
                                     job_id: int, task_id: str, datasource_id: str) -> AsyncIterator[ContentComponent]:
        """Extract text from image files using OCR - Non-blocking version"""
        
        try:
            # Run the blocking extraction in a thread pool
            loop = asyncio.get_running_loop()
            components = await loop.run_in_executor(
                None, 
                self._blocking_image_extraction,
                file_path, object_id, job_id, task_id, datasource_id
            )
            
            # Yield each component
            for component in components:
                yield component
                
        except Exception as e:
            yield self._create_error_component(object_id, f"Image file processing system error: {str(e)}")

    def _blocking_image_extraction(self, file_path: str, object_id: str, 
                                  job_id: int, task_id: str, datasource_id: str) -> List[ContentComponent]:
        """Synchronous image file extraction to run in thread pool"""
        components = []
        
        if not self._is_ocr_enabled():
            component = ContentComponent(
                object_id=object_id,
                component_type="no_content_extractable",
                component_id=f"{self._get_base_name(file_path)}_image_no_ocr_1",
                parent_path=file_path,
                content="OCR not available for image text extraction",
                original_size=self._get_file_size_sync(file_path),
                extracted_size=0,
                is_truncated=False,
                schema={},
                metadata=self._get_standard_metadata("no_content_extractable", {"reason": "ocr_disabled"}),
                extraction_method="ocr_check"
            )
            return [component]
        
        try:
            # Use synchronous OCR since we're already in thread pool
            ocr_text = self._extract_text_from_image(file_path)
            base_name = self._get_base_name(file_path)
            
            if ocr_text.strip():
                component = ContentComponent(
                    object_id=object_id,
                    component_type="image_ocr",
                    component_id=f"{base_name}_image_ocr_1",
                    parent_path=file_path,
                    content=ocr_text,
                    original_size=self._get_file_size_sync(file_path),
                    extracted_size=len(ocr_text),
                    is_truncated=False,
                    schema={},
                    metadata=self._get_standard_metadata("image_ocr", {"ocr_confidence": "filtered_0.5"}),
                    extraction_method="easyocr"
                )
                components.append(component)
            else:
                component = ContentComponent(
                    object_id=object_id,
                    component_type="no_content_extractable",
                    component_id=f"{base_name}_image_no_text_1",
                    parent_path=file_path,
                    content="No text detected in image",
                    original_size=self._get_file_size_sync(file_path),
                    extracted_size=0,
                    is_truncated=False,
                    schema={},
                    metadata=self._get_standard_metadata("no_content_extractable", {"reason": "no_text_detected"}),
                    extraction_method="easyocr"
                )
                components.append(component)
            
            return components
                
        except Exception as e:
            error_component = self._create_error_component(object_id, f"Image OCR failed: {str(e)}")
            return [error_component]

    def _get_file_size_sync(self, file_path: str) -> int:
        """Get file size synchronously for use in thread pool"""
        try:
            import os
            return os.path.getsize(file_path)
        except Exception:
            return 0

    # =============================================================================
    # Component Creation Helpers
    # =============================================================================

    def _create_text_component(self, object_id: str, component_id: str, file_path: str,
                             content: str, extra_metadata: Dict, extraction_method: str) -> ContentComponent:
        """Component size validation and standardized creation"""
        
        # Apply component size limits
        content_size_mb = len(content.encode('utf-8')) / (1024 * 1024)
        max_size_mb = self.extraction_config.limits.max_component_size_mb
        
        if content_size_mb > max_size_mb:
            content = self._apply_sampling_strategy(content, max_size_mb * 1024 * 1024)
            is_truncated = True
        else:
            # Apply text character limits
            content = content[:self.extraction_config.limits.max_text_chars]
            is_truncated = len(content) > self.extraction_config.limits.max_text_chars
        
        return ContentComponent(
            object_id=object_id,
            component_type="text",
            component_id=component_id,
            parent_path=file_path,
            content=content,
            original_size=len(content),
            extracted_size=len(content),
            is_truncated=is_truncated,
            schema={},
            metadata=self._get_standard_metadata("text", extra_metadata),
            extraction_method=extraction_method
        )

    def _get_standard_metadata(self, component_type: str, extra_metadata: Dict = None) -> Dict[str, Any]:
        """Standardized metadata structure"""
        
        standard_metadata = {
            "component_type": component_type,
            "extraction_timestamp": self._get_timestamp(),
            "extraction_config_version": "1.0"
        }
        
        if extra_metadata:
            standard_metadata.update(extra_metadata)
        
        return standard_metadata

    def _apply_sampling_strategy(self, content: str, max_bytes: int) -> str:
        """Implement sampling strategy for large content"""
        
        strategy = self.extraction_config.limits.sampling_strategy
        
        if strategy == "head":
            return content[:max_bytes]
        elif strategy == "tail":
            return content[-max_bytes:]
        elif strategy == "random":
            # Simple random sampling - take chunks from beginning, middle, end
            third = max_bytes // 3
            start = content[:third]
            middle_start = len(content) // 2 - third // 2
            middle = content[middle_start:middle_start + third]
            end = content[-third:]
            return start + "\n...\n" + middle + "\n...\n" + end
        else:
            return content[:max_bytes]  # Default to head

    def _get_timestamp(self) -> str:
        """Get current timestamp for metadata"""
        return datetime.now(timezone.utc).isoformat()

    # =============================================================================
    # Helper Methods - Fixed Component IDs and Added Missing Methods
    # =============================================================================

    def _get_base_name(self, file_path: str) -> str:
        """Get base filename without extension for component ID generation"""
        filename = os.path.basename(file_path)
        return filename.split('.')[0] if '.' in filename else filename

    def _create_error_component(self, object_id: str, error_message: str) -> ContentComponent:
        """Standardized error component creation"""
        base_name = self._get_base_name(object_id) if '/' in object_id else object_id
        return ContentComponent(
            object_id=object_id,
            component_type="extraction_error",
            component_id=f"{base_name}_error_1",
            parent_path="",
            content=error_message,
            original_size=0,
            extracted_size=len(error_message),
            is_truncated=False,
            schema={},
            metadata=self._get_standard_metadata("extraction_error", {"error": error_message}),
            extraction_method="error_handler"
        )

    def _create_oversized_component(self, object_id: str, file_info: Dict) -> ContentComponent:
        """Standardized oversized component creation"""
        base_name = self._get_base_name(file_info.get("path", object_id))
        return ContentComponent(
            object_id=object_id,
            component_type="file_too_large",
            component_id=f"{base_name}_oversized_1",
            parent_path=file_info.get("path", ""),
            content=f"File size {file_info.get('size_mb', 0):.2f}MB exceeds limit {self.extraction_config.limits.max_file_size_mb}MB",
            original_size=file_info.get("size_bytes", 0),
            extracted_size=0,
            is_truncated=False,
            schema={},
            metadata=self._get_standard_metadata("file_too_large", {"file_size_mb": file_info.get("size_mb", 0)}),
            extraction_method="size_check"
        )

    def _create_unsupported_component(self, object_id: str, file_path: str, file_type: str) -> ContentComponent:
        """Create component for unsupported file types"""
        base_name = self._get_base_name(file_path)
        return ContentComponent(
            object_id=object_id,
            component_type="unsupported_format",
            component_id=f"{base_name}_unsupported_1",
            parent_path=file_path,
            content=f"File type '{file_type}' not supported for extraction",
            original_size=os.path.getsize(file_path),
            extracted_size=0,
            is_truncated=False,
            schema={},
            metadata=self._get_standard_metadata("unsupported_format", {"file_type": file_type}),
            extraction_method="type_detection"
        )

    def _create_depth_exceeded_component(self, object_id: str, current_depth: int) -> ContentComponent:
        """Create component for depth exceeded errors"""
        base_name = self._get_base_name(object_id) if '/' in object_id else object_id
        return ContentComponent(
            object_id=object_id,
            component_type="extraction_error",
            component_id=f"{base_name}_depth_exceeded_1",
            parent_path="",
            content=f"Archive depth {current_depth} exceeds limit {self.extraction_config.limits.max_archive_depth}",
            original_size=0,
            extracted_size=0,
            is_truncated=False,
            schema={},
            metadata=self._get_standard_metadata("extraction_error", {"depth": current_depth}),
            extraction_method="depth_validation"
        )

    def _create_circular_reference_component(self, object_id: str, archive_path: str) -> ContentComponent:
        """Create component for circular reference errors"""
        base_name = self._get_base_name(archive_path)
        return ContentComponent(
            object_id=object_id,
            component_type="extraction_error",
            component_id=f"{base_name}_circular_ref_1",
            parent_path=archive_path,
            content=f"Circular archive reference detected: {archive_path}",
            original_size=0,
            extracted_size=0,
            is_truncated=False,
            schema={},
            metadata=self._get_standard_metadata("extraction_error", {"archive_path": archive_path}),
            extraction_method="circular_detection"
        )

    def _create_temp_path(self, job_id: int, task_id: str, datasource_id: str, suffix: str) -> str:
        """Create temp path with context IDs for uniqueness"""
        temp_base = getattr(self.system_config.system, 'temp_extraction_directory', tempfile.gettempdir())
        unique_suffix = uuid.uuid4().hex[:8]
        filename = f"extract_{job_id}_{task_id}_{datasource_id}_{suffix}_{unique_suffix}"
        return os.path.join(temp_base, filename)

    # =============================================================================
    # Configuration and OCR - Fixed OCR Configuration Path
    # =============================================================================

    def _is_ocr_enabled(self) -> bool:
        """Correct OCR configuration check"""
        # Fixed: Use correct system config path
        ocr_system_enabled = getattr(self.system_config, 'connector', None) and \
                           getattr(self.system_config.connector, 'easyocr_support_enabled', False)
        
        return (ocr_system_enabled and 
                self.extraction_config.features.ocr_enabled and
                self.ocr_available)

    def _setup_ocr_reader(self):
        """Initialize OCR reader with corrected configuration check"""
        # Check if OCR is enabled at system level
        ocr_system_enabled = getattr(self.system_config, 'connector', None) and \
                           getattr(self.system_config.connector, 'easyocr_support_enabled', False)
        
        if not ocr_system_enabled:
            self.logger.info("OCR disabled by system configuration")
            return
        
        try:
            import torch
            import easyocr
            
            use_gpu = torch.cuda.is_available()
            self.logger.info("Initializing OCR", gpu_available=use_gpu)
            
            self.ocr_reader = easyocr.Reader(['en'], gpu=use_gpu)
            self.ocr_available = True
            
        except ImportError as e:
            self.logger.warning("EasyOCR not available", error=str(e))
            self.ocr_available = False
            self.ocr_reader = None


    async def _extract_text_from_image_async(self, image_path: str) -> str:
        """Extract text from image using OCR - truly async version"""
        
        if not self.ocr_available:
            return ""
        
        try:
            # Run the blocking OCR operation in a thread pool
            loop = asyncio.get_running_loop()
            extracted_text = await loop.run_in_executor(
                None, 
                self._blocking_ocr_extraction,
                image_path
            )
            return extracted_text
            
        except Exception as e:
            self.logger.warning("OCR extraction failed", 
                               image_path=image_path, error=str(e))
            return ""

    def _blocking_ocr_extraction(self, image_path: str) -> str:
        """Synchronous OCR extraction to run in thread pool"""
        
        try:
            # EasyOCR readtext is CPU-intensive and blocking
            results = self.ocr_reader.readtext(image_path)
            extracted_text = " ".join([
                result[1] for result in results 
                if result[2] > 0.5  # Confidence threshold
            ])
            return extracted_text
            
        except Exception as e:
            self.logger.warning("OCR processing failed", 
                               image_path=image_path, error=str(e))
            return ""

    def _extract_text_from_image(self, image_path: str) -> str:
        """Extract text from image using OCR - sync version for compatibility"""
        if not self.ocr_available:
            return ""
        
        try:
            results = self.ocr_reader.readtext(image_path)
            extracted_text = " ".join([
                result[1] for result in results 
                if result[2] > 0.5  # Confidence threshold
            ])
            return extracted_text
            
        except Exception as e:
            self.logger.warning("OCR extraction failed", 
                               image_path=image_path, error=str(e))
            return ""

    # =============================================================================
    # Table Conversion Methods - All Fixed for Empty Content
    # =============================================================================

    def _convert_pymupdf_table(self, df, table_index: int) -> Dict[str, Any]:
        """Convert PyMuPDF pandas DataFrame to structured format"""
        headers = df.columns.tolist()
        rows = []
        
        for _, row in df.iterrows():
            row_dict = {headers[i]: str(row.iloc[i]) for i in range(len(headers))}
            rows.append(row_dict)
        
        return {
            "headers": headers,
            "rows": rows,
            "row_count": len(rows),
            "column_count": len(headers),
            "serialized_rows": "\n".join(["|".join(row_dict.values()) for row_dict in rows]),
            "table_index": table_index
        }

    def _convert_camelot_table(self, camelot_table) -> Dict[str, Any]:
        """Convert Camelot table to structured format"""
        df = camelot_table.df
        headers = df.columns.tolist()
        rows = []
        
        for _, row in df.iterrows():
            row_dict = {headers[i]: str(row.iloc[i]) for i in range(len(headers))}
            rows.append(row_dict)
        
        return {
            "headers": headers,
            "rows": rows,
            "row_count": len(rows),
            "column_count": len(headers),
            "serialized_rows": "\n".join(["|".join(row_dict.values()) for row_dict in rows])
        }

    def _convert_tabula_table(self, tabula_df) -> Dict[str, Any]:
        """Convert Tabula DataFrame to structured format"""
        headers = tabula_df.columns.tolist()
        rows = []
        
        for _, row in tabula_df.iterrows():
            row_dict = {headers[i]: str(row.iloc[i]) for i in range(len(headers))}
            rows.append(row_dict)
        
        return {
            "headers": headers,
            "rows": rows,
            "row_count": len(rows),
            "column_count": len(headers),
            "serialized_rows": "\n".join(["|".join(row_dict.values()) for row_dict in rows])
        }

    def _convert_docx_table(self, docx_table) -> Dict[str, Any]:
        """Convert Word document table to structured format"""
        headers = []
        rows = []
        
        if docx_table.rows:
            first_row = docx_table.rows[0]
            headers = [cell.text.strip() for cell in first_row.cells]
        
        for row in docx_table.rows[1:]:
            row_data = {}
            for i, cell in enumerate(row.cells):
                header = headers[i] if i < len(headers) else f"column_{i}"
                row_data[header] = cell.text.strip()
            rows.append(row_data)
        
        return {
            "headers": headers,
            "rows": rows,
            "row_count": len(rows),
            "column_count": len(headers),
            "serialized_rows": "\n".join(["|".join(row_data.values()) for row_data in rows])
        }

    def _convert_pptx_table(self, pptx_table) -> Dict[str, Any]:
        """Convert PowerPoint table to structured format"""
        headers = []
        rows = []
        
        # Extract headers from first row
        if pptx_table.rows:
            first_row = pptx_table.rows[0]
            headers = [cell.text.strip() for cell in first_row.cells]
        
        # Extract data rows
        for row in pptx_table.rows[1:]:
            row_data = {}
            for i, cell in enumerate(row.cells):
                header = headers[i] if i < len(headers) else f"column_{i}"
                row_data[header] = cell.text.strip()
            rows.append(row_data)
        
        return {
            "headers": headers,
            "rows": rows,
            "row_count": len(rows),
            "column_count": len(headers),
            "serialized_rows": "\n".join(["|".join(row_data.values()) for row_data in rows])
        }

    def _convert_excel_sheet(self, df, sheet_name: str) -> Dict[str, Any]:
        """Convert pandas DataFrame to structured format"""
        headers = df.columns.tolist()
        rows = []
        
        for _, row in df.iterrows():
            row_dict = {headers[i]: str(row.iloc[i]) for i in range(len(headers))}
            rows.append(row_dict)
        
        return {
            "headers": headers,
            "rows": rows,
            "row_count": len(rows),
            "column_count": len(headers),
            "serialized_rows": "\n".join(["|".join(row_dict.values()) for row_dict in rows]),
            "sheet_name": sheet_name
        }

    def _convert_html_table(self, table_element) -> Dict[str, Any]:
        """Convert BeautifulSoup table to structured format"""
        headers = []
        rows = []
        
        # Extract headers
        header_row = table_element.find('tr')
        if header_row:
            headers = [th.text.strip() for th in header_row.find_all(['th', 'td'])]
        
        # Extract data rows
        data_rows = table_element.find_all('tr')[1:] if headers else table_element.find_all('tr')
        for row in data_rows:
            cells = [td.text.strip() for td in row.find_all(['td', 'th'])]
            if cells and len(cells) <= len(headers):
                # Pad cells if fewer than headers
                while len(cells) < len(headers):
                    cells.append("")
                row_data = {headers[i]: cells[i] for i in range(len(headers))}
                rows.append(row_data)
        
        return {
            "headers": headers,
            "rows": rows,
            "row_count": len(rows),
            "column_count": len(headers),
            "serialized_rows": "\n".join(["|".join(row_data.values()) for row_data in rows])
        }

    # =============================================================================
    # File Type Detection
    # =============================================================================

    def _detect_file_type(self, file_path: str) -> str:
        """Detect file type using MIME type detection with extension fallback"""
        try:
            try:
                import magic
                mime_type = magic.from_file(file_path, mime=True)
                
                mime_to_type = {
                    'application/pdf': 'pdf',
                    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
                    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
                    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
                    'text/html': 'html',
                    'text/plain': 'txt',
                    'application/json': 'json',
                    'application/xml': 'xml',
                    'text/xml': 'xml',
                    'application/zip': 'zip',
                    'application/x-tar': 'tar',
                    'application/gzip': 'tar',
                    'image/png': 'image',
                    'image/jpeg': 'image',
                    'image/tiff': 'image',
                    'message/rfc822': 'eml'
                }
                
                if mime_type in mime_to_type:
                    return mime_to_type[mime_type]
                    
            except ImportError:
                self.logger.debug("python-magic not available, using extension fallback")
            
            return self._detect_type_by_extension(file_path)
            
        except Exception as e:
            self.logger.warning("File type detection failed", file_path=file_path, error=str(e))
            return 'unknown'

    def _detect_type_by_extension(self, file_path: str) -> str:
        """Fallback file type detection by extension"""
    # Check compound extensions first
        file_path_lower = file_path.lower()
        if file_path_lower.endswith('.tar.gz'):
            return 'tar'
        if file_path_lower.endswith('.tar.bz2'):
            return 'tar'        
        ext = Path(file_path).suffix.lower()
        
        extension_map = {
            '.pdf': 'pdf',
            '.docx': 'docx',
            '.xlsx': 'xlsx', 
            '.pptx': 'pptx',
            '.html': 'html',
            '.htm': 'html',
            '.txt': 'txt',
            '.json': 'json',
            '.xml': 'xml',
            '.zip': 'zip',
            '.tar': 'tar',
            
            '.tgz': 'tar',
            '.png': 'image',
            '.jpg': 'image',
            '.jpeg': 'image',
            '.tiff': 'image',
            '.eml': 'eml'
        }
        
        return extension_map.get(ext, 'unsupported')

    # =============================================================================
    # Security and Validation
    # =============================================================================

    def _validate_archive_member_path(self, member_path: str) -> bool:
        """Validate archive member path for security"""
        
        # Check for absolute paths
        if member_path.startswith('/') or member_path.startswith('\\'):
            return False
        
        # Check for parent directory references
        path_parts = member_path.replace('\\', '/').split('/')
        if '..' in path_parts:
            return False
        
        # Check for Windows drive letters
        if ':' in member_path and len(member_path) > 1 and member_path[1] == ':':
            return False
        
        return True

    # =============================================================================
    # Resource Management
    # =============================================================================

    async def _cleanup_temp_files_async(self, extraction_id: str):
        """Clean up temporary files created during extraction - async version"""
        
        cleanup_success = []
        cleanup_failures = []
        
        for temp_file in self.temp_files_created:
            try:
                if await aiofiles.os.path.exists(temp_file):
                    await aiofiles.os.remove(temp_file)
                    cleanup_success.append(temp_file)
                
            except Exception as cleanup_error:
                cleanup_failures.append({"file": temp_file, "error": str(cleanup_error)})
                self.logger.warning("Temp file cleanup failed",
                                   temp_file=temp_file,
                                   error=str(cleanup_error))
        
        self.temp_files_created.clear()
        
        if cleanup_success or cleanup_failures:
            self.logger.info("Resource cleanup completed",
                             extraction_id=extraction_id,
                             files_cleaned=len(cleanup_success),
                             cleanup_failures=len(cleanup_failures))

    def _cleanup_temp_files(self, extraction_id: str):
        """Clean up temporary files created during extraction - sync version for compatibility"""
        
        cleanup_success = []
        cleanup_failures = []
        
        for temp_file in self.temp_files_created:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
                    cleanup_success.append(temp_file)
                
            except Exception as cleanup_error:
                cleanup_failures.append({"file": temp_file, "error": str(cleanup_error)})
                self.logger.warning("Temp file cleanup failed",
                                   temp_file=temp_file,
                                   error=str(cleanup_error))
        
        self.temp_files_created.clear()
        
        if cleanup_success or cleanup_failures:
            self.logger.info("Resource cleanup completed",
                             extraction_id=extraction_id,
                             files_cleaned=len(cleanup_success),
                             cleanup_failures=len(cleanup_failures))

    # =============================================================================
    # Library Registry and Configuration  
    # =============================================================================

    def _build_extractor_registry(self) -> Dict[str, callable]:
        """Build registry of available extractors based on library availability"""
        
        registry = {}
        missing_libraries = []
        
        # Always available
        registry["txt"] = self.extract_txt_components
        registry["json"] = self.extract_json_components
        registry["xml"] = self.extract_xml_components
        registry["eml"] = self.extract_eml_components
        
        # PDF processing
        try:
            registry["pdf"] = self.extract_pdf_components
        except ImportError:
            missing_libraries.append("pymupdf")
        
        # Office processing
        try:
            registry["docx"] = self.extract_docx_components
        except ImportError:
            missing_libraries.append("python-docx")
        
        try:
            registry["xlsx"] = self.extract_xlsx_components
        except ImportError:
            missing_libraries.append("pandas/openpyxl")
        
        try:
            registry["pptx"] = self.extract_pptx_components
        except ImportError:
            missing_libraries.append("python-pptx")
        
        # HTML processing
        try:
            registry["html"] = self.extract_html_components
        except ImportError:
            missing_libraries.append("beautifulsoup4")
        
        # Image processing
        registry["image"] = self.extract_image_components
        
        # Note: Archives handled separately to avoid recursion
        
        self.logger.info("Extractor registry built",
                        available_extractors=list(registry.keys()),
                        missing_libraries=missing_libraries)
        
        return registry


    def _validate_extraction_config(self):
            """Validate configuration and log issues"""
            
            validation_errors = []
            
            # OCR flag validation
            ocr_system_enabled = getattr(self.system_config, 'connector', None) and \
                               getattr(self.system_config.connector, 'easyocr_support_enabled', False)
            
            if (not ocr_system_enabled and 
                self.extraction_config.features.ocr_enabled):
                validation_errors.append({
                    "type": "config_conflict",
                    "message": "OCR enabled in datasource but disabled at system level",
                    "resolution": "Using system setting (OCR disabled)"
                })
            
            # Resource limit validation
            total_memory_limit = self.system_config.system.total_process_memory_limit_mb
            safe_memory_limit = total_memory_limit * 0.7
            if self.extraction_config.limits.max_file_size_mb > safe_memory_limit:
                validation_errors.append({
                    "type": "resource_limit_exceeded",
                    "message": f"max_file_size_mb exceeds safe memory limit",
                    "resolution": "File size limit should be reduced"
                })
            
            if validation_errors:
                self.logger.warning("Configuration validation issues found",
                                   validation_errors=validation_errors)
                             
    # =============================================================================
    # Simple File Type Extractors - Signatures Updated for Context IDs
    # =============================================================================

    async def extract_txt_components(self, file_path: str, object_id: str,
                                   job_id: int, task_id: str, datasource_id: str) -> AsyncIterator[ContentComponent]:
        """Extract text file content"""
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = await f.read()
            
            text_component = self._create_text_component(
                object_id, f"{self._get_base_name(file_path)}_txt_text_1", file_path,
                text_content, {"encoding": "utf-8"}, "text_file_read"
            )
            yield text_component
            
        except Exception as e:
            yield self._create_error_component(object_id, f"Text file extraction failed: {str(e)}")

    async def extract_json_components(self, file_path: str, object_id: str,
                                    job_id: int, task_id: str, datasource_id: str) -> AsyncIterator[ContentComponent]:
        """Extract JSON file content"""
        try:
            import json
            
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                json_text = await f.read()
                json_data = json.loads(json_text)
            
            formatted_json = json.dumps(json_data, indent=2, ensure_ascii=False)
            
            text_component = self._create_text_component(
                object_id, f"{self._get_base_name(file_path)}_json_text_1", file_path,
                formatted_json, {"json_type": type(json_data).__name__}, "json_parsing"
            )
            yield text_component
            
        except Exception as e:
            yield self._create_error_component(object_id, f"JSON extraction failed: {str(e)}")

    async def extract_image_components(self, file_path: str, object_id: str,
                                     job_id: int, task_id: str, datasource_id: str) -> AsyncIterator[ContentComponent]:
        """Extract text from image files using OCR"""
        
        if not self._is_ocr_enabled():
            yield ContentComponent(
                object_id=object_id,
                component_type="no_content_extractable",
                component_id=f"{self._get_base_name(file_path)}_image_no_ocr_1",
                parent_path=file_path,
                content="OCR not available for image text extraction",
                original_size=(await aiofiles.os.stat(file_path)).st_size,
                extracted_size=0,
                is_truncated=False,
                schema={},
                metadata=self._get_standard_metadata("no_content_extractable", {"reason": "ocr_disabled"}),
                extraction_method="ocr_check"
            )
            return
        
        try:
            ocr_text = await self._extract_text_from_image_async(file_path)
            base_name = self._get_base_name(file_path)
            
            if ocr_text.strip():
                yield ContentComponent(
                    object_id=object_id,
                    component_type="image_ocr",
                    component_id=f"{base_name}_image_ocr_1",
                    parent_path=file_path,
                    content=ocr_text,
                    original_size=(await aiofiles.os.stat(file_path)).st_size,
                    extracted_size=len(ocr_text),
                    is_truncated=False,
                    schema={},
                    metadata=self._get_standard_metadata("image_ocr", {"ocr_confidence": "filtered_0.5"}),
                    extraction_method="easyocr"
                )
            else:
                yield ContentComponent(
                    object_id=object_id,
                    component_type="no_content_extractable",
                    component_id=f"{base_name}_image_no_text_1",
                    parent_path=file_path,
                    content="No text detected in image",
                    original_size=(await aiofiles.os.stat(file_path)).st_size,
                    extracted_size=0,
                    is_truncated=False,
                    schema={},
                    metadata=self._get_standard_metadata("no_content_extractable", {"reason": "no_text_detected"}),
                    extraction_method="easyocr"
                )
                
        except Exception as e:
            yield self._create_error_component(object_id, f"Image OCR failed: {str(e)}")